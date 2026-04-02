"""PPTX Engine — PowerPoint generation via python-pptx."""

from pathlib import Path

from generate import register_engine


@register_engine('pptx')
class PptxEngine:
    def generate(self, data, template_path, output_path, schema=None):
        """
        Generate PPTX from data.

        Expected data format:
        {
            "title": "Presentation Title",
            "subtitle": "Optional subtitle",
            "author": "Author Name",
            "slides": [
                {"title": "Slide Title", "content": "Bullet text or paragraph", "layout": "content"},
                {"title": "Data Slide", "items": [{"label": "X", "value": 10}], "layout": "table"}
            ]
        }
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise RuntimeError("python-pptx required: pip install python-pptx")

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

        # Use template if provided
        if template_path and Path(template_path).exists() and str(template_path).endswith('.pptx'):
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()

        # Brand colors
        brand_dark = RGBColor(0x16, 0x21, 0x3E)
        brand_white = RGBColor(0xFF, 0xFF, 0xFF)

        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = data.get('title', 'Presentation')
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = data.get('subtitle', data.get('date', ''))

        # Content slides
        for slide_data in data.get('slides', []):
            layout_type = slide_data.get('layout', 'content')

            if layout_type == 'section':
                # Section header
                slide_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get('title', '')

            elif layout_type == 'table' and 'items' in slide_data:
                # Table slide
                slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get('title', '')

                items = slide_data['items']
                if items:
                    rows = len(items) + 1
                    cols = len(items[0])
                    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(0.5 * rows)).table

                    # Headers
                    for col_idx, key in enumerate(items[0].keys()):
                        table.cell(0, col_idx).text = str(key)

                    # Data
                    for row_idx, item in enumerate(items, 1):
                        for col_idx, val in enumerate(item.values()):
                            table.cell(row_idx, col_idx).text = str(val)

            else:
                # Content slide (bullets or paragraph)
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get('title', '')

                body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                if body:
                    content = slide_data.get('content', '')
                    if isinstance(content, list):
                        tf = body.text_frame
                        tf.text = content[0] if content else ''
                        for bullet in content[1:]:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    else:
                        body.text = content

        prs.save(str(output_path))
